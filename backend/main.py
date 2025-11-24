import os
import json
import uuid
from datetime import datetime, timedelta
from pathlib import Path
from collections import defaultdict
from typing import List, Dict, Any, Tuple
from io import BytesIO
import bcrypt
import pdfplumber
from pdf2image import convert_from_bytes
import pytesseract
from PIL import Image
from bs4 import BeautifulSoup
import pandas as pd
import requests
import asyncio
from google.oauth2 import service_account
import google.auth.transport.requests
from fastapi import Query, Header, FastAPI, UploadFile, File, Form, Request, Depends, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from pydantic import BaseModel
from dotenv import load_dotenv
from rank_bm25 import BM25Okapi
import re
from sentence_transformers import CrossEncoder
reranker = CrossEncoder('BAAI/bge-reranker-base')  # or 'bge-reranker-large', or Cohere's rerank model if paying
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
load_dotenv()

# ----------- OPENROUTER UTILITY -----------
def call_openrouter_model(model_name, prompt, api_key=None, temperature=0.7, max_tokens=1200):
    api_key = api_key or os.getenv("OPENROUTER_API_KEY2")
    url = "https://openrouter.ai/api/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
        "HTTP-Referer": "http://localhost:3000",
        "X-Title": "ASKEASE"
    }
    payload = {
        "model": model_name,
        "messages": [
            {"role": "system", "content": "You are a helpful assistant for academic and technical topics."},
            {"role": "user", "content": prompt}
        ],
        "temperature": temperature,
        "max_tokens": max_tokens
    }
    resp = requests.post(url, headers=headers, json=payload, timeout=60)
    try:
        data = resp.json()
        return data["choices"][0]["message"]["content"]
    except Exception as e:
        print("OpenRouter parse error:", e, resp.text)
        return "Error: Could not parse OpenRouter response."

# ---------- SQLAlchemy & MySQL Setup ----------
from sqlalchemy import create_engine, Column, String, Text, DateTime, JSON, ForeignKey
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, relationship, Session

DATABASE_URL = os.getenv("DATABASE_URL")
if not DATABASE_URL:
    raise RuntimeError("Please set DATABASE_URL in .env")

Base = declarative_base()

class User(Base):
    __tablename__ = "users"
    id = Column(String(36), primary_key=True, default=lambda: str(uuid.uuid4()))
    username = Column(String(50), unique=True, nullable=False)
    password_hash = Column(String(64), nullable=False)
    email = Column(String(100), nullable=False)
    created_at = Column(DateTime, default=datetime.now)
    updated_at = Column(DateTime, default=datetime.now, onupdate=datetime.now)
    sessions = relationship("ChatSession", back_populates="user")

class UserSession(Base):
    __tablename__ = "user_sessions"
    id = Column(String(36), primary_key=True, default=lambda: str(uuid.uuid4()))
    user_id = Column(String(36), ForeignKey("users.id"), nullable=False)
    session_token = Column(String(200), unique=True, nullable=False)
    created_at = Column(DateTime, default=datetime.now)
    expires_at = Column(DateTime, nullable=False)

class ChatSession(Base):
    __tablename__ = "chat_sessions"
    id = Column(String(36), primary_key=True, default=lambda: str(uuid.uuid4()))
    user_id = Column(String(36), ForeignKey("users.id"), nullable=True)
    title = Column(String(200), nullable=True)
    created_at = Column(DateTime, default=datetime.now)
    updated_at = Column(DateTime, default=datetime.now, onupdate=datetime.now)
    user = relationship("User", back_populates="sessions")
    messages = relationship("ChatMessage", back_populates="session")

class ChatMessage(Base):
    __tablename__ = "chat_messages"
    id = Column(String(36), primary_key=True, default=lambda: str(uuid.uuid4()))
    session_id = Column(String(36), ForeignKey("chat_sessions.id"), nullable=False)
    sender = Column(String(20), nullable=False)
    message_text = Column(Text, nullable=False)
    citations = Column(JSON, nullable=True)
    created_at = Column(DateTime, default=datetime.now)
    session = relationship("ChatSession", back_populates="messages")

engine = create_engine(DATABASE_URL, future=True)
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base.metadata.create_all(bind=engine)

def get_db():
    db = SessionLocal()
    try: yield db
    finally: db.close()

# ---------- EMBEDDINGS & LLM ----------
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import RetrievalQA

STORAGE_ROOT = Path(os.getenv("STORAGE_DIR","storage"))
STORAGE_ROOT.mkdir(parents=True, exist_ok=True)
embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

prompt_template = """You are a precise data-extraction assistant. Use ONLY the CONTEXT below (do NOT invent facts).
Follow these rules exactly:
1) Answer using only text from CONTEXT. If information is not present, reply exactly: Not found
2) If the CONTEXT contains a table in Markdown, copy and return the table as Markdown (don't restate rows/columns—show the table).
3) Preserve numbers, units, and currency exactly as in CONTEXT.
4) If asked for multiple matches, return all matches.
5) Provide short answers (concise), then add a "Sources" section listing source filenames and locations.
6) If you extract a table, present it in Markdown format only.
7) If asked to summarize, keep summary < 120 words and do not invent facts.
Context:
{context}
Question: {question}
Answer:"""

from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation

def extract_text_from_file_upload(upload_file: UploadFile) -> Tuple[List[str], List[Dict]]:
    import re
    import pdfplumber
    import pandas as pd
    from docx import Document as DocxDocument
    from pptx import Presentation as PptxPresentation
    from PIL import Image
    from bs4 import BeautifulSoup
    from io import BytesIO

    filename = upload_file.filename
    ext = Path(filename).suffix.lower()

    # FIX: Always seek(0) before read, only read ONCE
    try:
        upload_file.file.seek(0)
    except Exception:
        pass
    file_bytes = upload_file.file.read()
    texts, metas = [], []

    def add(t, meta):
        if t and t.strip():
            texts.append(t)
            metas.append(meta)

    def markdown_table(table_data):
        if not table_data or not all(isinstance(row, (list, tuple)) for row in table_data): 
            return ""
        header = "| " + " | ".join(table_data[0]) + " |"
        sep = "| " + " | ".join(['---'] * len(table_data[0])) + " |"
        body = ["| " + " | ".join(row) + " |" for row in table_data[1:]]
        return "\n".join([header, sep] + body)

    def extract_code_blocks(text):
        blocks = []
        blocks += re.findall(r"``````", text, re.DOTALL)
        blocks += re.findall(r'(import\s+\w+[\s\S]+?\})', text)
        blocks += re.findall(r'(public\s+class[\s\S]+?\})', text)
        blocks += re.findall(r'(def\s+\w+[\s\S]+?):\n(?: {4,}.*\n?)+', text)
        code_chunks = []
        for cb in blocks:
            cb2 = cb.strip()
            if cb2 and cb2 not in code_chunks:
                code_chunks.append(cb2)
        return code_chunks

    def extract_text_with_ocr(pdf_bytes, page_num):
        """Extract text from PDF page using OCR if regular extraction fails"""
        try:
            # Convert the specific page to image - fix path formatting
            images = convert_from_bytes(
                pdf_bytes, 
                first_page=page_num+1, 
                last_page=page_num+1, 
                poppler_path=r"C:\poppler\poppler-25.07.0\Library\bin"  # Use raw string
            )
            if images:
                # Perform OCR on the image
                ocr_text = pytesseract.image_to_string(images[0])
                return ocr_text.strip()
        except Exception as e:
            print(f"OCR failed for page {page_num}: {e}")
        return ""

    try:
        if ext == ".pdf":
            # Use ONLY BytesIO(file_bytes)
            with pdfplumber.open(BytesIO(file_bytes)) as pdf:
                for i, page in enumerate(pdf.pages):
                    # First try regular text extraction
                    t = page.extract_text() or ""
                    
                    # If no text found, try OCR
                    if not t.strip():
                        print(f"Page {i+1}: No text found, attempting OCR...")
                        t = extract_text_with_ocr(file_bytes, i)
                        if t.strip():
                            print(f"Page {i+1}: OCR extracted {len(t)} characters")
                    
                    # Extract tables
                    tables = page.extract_tables()
                    code_chunks = extract_code_blocks(t)
                    
                    # Add code chunks
                    for code in code_chunks:
                        add(code, {"source": filename, "page": i+1, "type": "code"})
                    
                    # Add tables
                    for j, table in enumerate(tables):
                        md_table = markdown_table(table)
                        add(md_table, {"source": filename, "page": i+1, "type": "table"})
                    
                    # Add text content (whether from regular extraction or OCR)
                    if t.strip():
                        add(t, {"source": filename, "page": i+1, "type": "text"})
                    else:
                        print(f"Page {i+1}: No text content available even with OCR")

        elif ext == ".docx":
            doc = DocxDocument(BytesIO(file_bytes))
            full_txt = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
            code_chunks = extract_code_blocks(full_txt)
            for code in code_chunks:
                add(code, {"source": filename, "type": "code"})
            for ti, table in enumerate(doc.tables):
                table_data = [[cell.text for cell in row.cells] for row in table.rows]
                md_table = markdown_table(table_data)
                add(md_table, {"source": filename, "type": "table"})
            if full_txt.strip():
                add(full_txt, {"source": filename, "type": "text"})

        elif ext in [".txt", ".md"]:
            body = file_bytes.decode("utf-8", errors="ignore")
            code_chunks = extract_code_blocks(body)
            for code in code_chunks:
                add(code, {"source": filename, "type": "code"})
            if body.strip():
                add(body, {"source": filename, "type": "text"})

        elif ext == ".csv":
            df = pd.read_csv(BytesIO(file_bytes))
            table_data = [list(df.columns)] + df.astype(str).values.tolist()
            md_table = markdown_table(table_data)
            add(md_table, {"source": filename, "type": "table"})

        elif ext in [".xlsx", ".xls"]:
            wb = pd.ExcelFile(BytesIO(file_bytes))
            for sheet_name in wb.sheet_names:
                df = wb.parse(sheet_name)
                table_data = [list(df.columns)] + df.astype(str).values.tolist()
                md_table = markdown_table(table_data)
                add(md_table, {"source": filename, "sheet": sheet_name, "type": "table"})

        elif ext == ".pptx":
            prs = PptxPresentation(BytesIO(file_bytes))
            for i, slide in enumerate(prs.slides):
                slide_lines = []
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and getattr(shape, "text_frame") is not None:
                        for para in shape.text_frame.paragraphs:
                            txt = para.text.strip()
                            if txt:
                                slide_lines.append(txt)
                    elif hasattr(shape, "text") and shape.text and shape.text.strip():
                        slide_lines.append(shape.text.strip())
                joined = "\n".join(slide_lines)
                print(f"\n=== SLIDE {i+1} ===\n{joined}\n")  # Debug/Remove in prod
                if joined.strip():
                    add(joined, {"source": filename, "slide": i+1, "type": "text"})
                # Also extract code blocks from joined slide text
                code_chunks = extract_code_blocks(joined)
                for code in code_chunks:
                    add(code, {"source": filename, "slide": i+1, "type": "code"})

        elif ext in [".html", ".htm"]:
            soup = BeautifulSoup(file_bytes.decode("utf-8", "ignore"), "html.parser")
            for tag in soup.find_all(["code", "pre"]):
                code_text = tag.get_text(separator="\n").strip()
                if code_text:
                    add(code_text, {"source": filename, "type": "code"})
            add(soup.get_text(separator="\n"), {"source": filename, "type": "text"})
            for t_index, table in enumerate(soup.find_all("table")):
                rows = []
                for tr in table.find_all("tr"):
                    rows.append([td.get_text(strip=True) for td in tr.find_all(["td", "th"])])
                md_table = markdown_table(rows)
                add(md_table, {"source": filename, "type": "table"})

        elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"]:
            img = Image.open(BytesIO(file_bytes))
            ocr_text = pytesseract.image_to_string(img)
            add(ocr_text, {"source": filename, "type": "ocr"})

        else:
            try:
                t = file_bytes.decode("utf-8", errors="ignore")
                for code in extract_code_blocks(t):
                    add(code, {"source": filename, "type": "code"})
                if t.strip():
                    add(t, {"source": filename, "type": "text"})
                else:
                    raise ValueError("Unsupported or binary file")
            except Exception:
                raise ValueError(f"Unsupported file type {ext}")
    
    except Exception as e:
        print(f"Error processing file {filename}: {e}")
        # Try to extract any text using OCR as last resort for PDFs
        if ext == ".pdf":
            print("Attempting fallback OCR for entire PDF...")
            try:
                images = convert_from_bytes(file_bytes)
                for i, image in enumerate(images):
                    ocr_text = pytesseract.image_to_string(image)
                    if ocr_text.strip():
                        add(ocr_text, {"source": filename, "page": i+1, "type": "ocr"})
                        print(f"Fallback OCR extracted text from page {i+1}")
            except Exception as ocr_error:
                print(f"Fallback OCR also failed: {ocr_error}")
        
        if not texts:  # If still no text, re-raise the original error
            raise e
    
    finally:
        try:
            upload_file.file.close()
        except Exception:
            pass
    
    print(f"Extracted {len(texts)} text chunks from {filename}")
    return texts, metas


class SimpleBM25Retriever:
    def __init__(self, docs, chunk_metadatas):
        self.docs = docs
        self.chunk_metadatas = chunk_metadatas
        self.corpus = [self._preprocess_text(doc) for doc in docs]
        self.bm25 = BM25Okapi([d.split() for d in self.corpus])

    def _preprocess_text(self, text):
        return re.sub(r"[^\w\s]", " ", text).lower()

    def get_relevant_documents(self, query, k=5):
        processed_query = self._preprocess_text(query)
        scores = self.bm25.get_scores(processed_query.split())
        topk = sorted(enumerate(scores), key=lambda x: -x[1])[:k]
        return [(self.docs[i], self.chunk_metadatas[i]) for i, _ in topk if scores[i] > 0]

def rerank_chunks(query, docs_with_metas, top_k=1000):
    """
    Reranks input doc objects with .metadata and .page_content using cross-encoder.
    Returns list of top_k doc objects (sorted by reranker score).
    """
    sentences = [getattr(d, "page_content", "") for d in docs_with_metas]
    pairs = [(query, s) for s in sentences]
    scores = reranker.predict(pairs)
    reranked = sorted(zip(scores, docs_with_metas), key=lambda x: -x[0])
    return [doc for score, doc in reranked[:top_k]]

app = FastAPI(title="RAG + OCR + OpenRouter fallback backend")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000","http://127.0.0.1:3000","*"],
    allow_credentials=True,allow_methods=["*"],allow_headers=["*"],
)
vectorstore_dict: Dict[str, Any] = {}
qa_chain_dict: Dict[str, Any] = {}
qa_history_dict: Dict[str, List[dict]] = {}

@app.post("/upload")
async def upload(
    request: Request,
    files: List[UploadFile] = File(...),
    db: Session = Depends(get_db),
    session_token: str = Form(...),
    session_id: str = Form(None)  # ← NEW: Accept existing session ID
):
    if not files:
        return JSONResponse({"error": "No files provided"}, status_code=400)
    if not session_token:
        return JSONResponse({"error": "Login required to upload files"}, status_code=401)
    
    user_session = db.query(UserSession).filter(UserSession.session_token == session_token).first()
    if not user_session or user_session.expires_at < datetime.now():
        return JSONResponse({"error": "Invalid or expired session token"}, status_code=401)
    
    user_id = user_session.user_id
    all_texts = []
    all_metas = []
    errors = []
    
    for f in files:
        try:
            texts, metas = extract_text_from_file_upload(f)
            if texts:
                all_texts.extend(texts)
                all_metas.extend(metas)
            else:
                errors.append(f"{f.filename}: no text extracted")
        except Exception as e:
            errors.append(f"{f.filename}: {str(e)}")
    
    if not all_texts:
        return JSONResponse({"error": "No text extracted", "details": errors}, status_code=400)

    # === NEW: Check if we should add to existing session ===
    if session_id and session_id.strip():
        # ADD TO EXISTING SESSION
        sid = session_id.strip()
        existing_session = db.query(ChatSession).filter(ChatSession.id == sid).first()
        if not existing_session:
            return JSONResponse({"error": "Invalid existing session ID"}, status_code=400)
        
        # Load existing vector store from disk
        storage_path = STORAGE_ROOT / sid
        if not storage_path.exists():
            return JSONResponse({"error": "Existing session data not found"}, status_code=400)
        
        try:
            # Load existing vector store
            existing_vectorstore = FAISS.load_local(str(storage_path), embeddings, allow_dangerous_deserialization=True)
            
            # Process new files
            splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
            new_docs_for_embeddings = []
            new_metadatas_for_embeddings = []
            
            for t, m in zip(all_texts, all_metas):
                pieces = splitter.split_text(t)
                for i, p in enumerate(pieces):
                    new_docs_for_embeddings.append(p)
                    md = m.copy()
                    md["chunk"] = i
                    new_metadatas_for_embeddings.append(md)
            
            # Add new documents to existing vector store
            existing_vectorstore.add_texts(new_docs_for_embeddings, metadatas=new_metadatas_for_embeddings)
            
            # Save updated vector store
            existing_vectorstore.save_local(str(storage_path))
            
            # Update in-memory stores
            vectorstore_dict[sid] = existing_vectorstore
            
            # Update BM25 retriever with new documents
            if sid in qa_chain_dict:
                # Get existing docs and add new ones
                existing_retriever = qa_chain_dict[sid]["bm25"]
                all_docs = existing_retriever.docs + new_docs_for_embeddings
                all_metadatas = existing_retriever.chunk_metadatas + new_metadatas_for_embeddings
                updated_retriever = SimpleBM25Retriever(all_docs, all_metadatas)
                qa_chain_dict[sid] = {"bm25": updated_retriever}
            else:
                # Create new BM25 retriever with all documents
                qa_chain_dict[sid] = {"bm25": SimpleBM25Retriever(new_docs_for_embeddings, new_metadatas_for_embeddings)}
            
            return {
                "session_id": sid,
                "message": f"✅ {len(files)} file(s) ADDED to existing session. Total chunks: {existing_vectorstore.index.ntotal}",
                "errors": errors,
                "action": "added_to_existing"
            }
            
        except Exception as e:
            print(f"Error adding to existing session: {e}")
            return JSONResponse({"error": f"Failed to add files to existing session: {str(e)}"}, status_code=500)
    
    else:
        # CREATE NEW SESSION (original behavior)
        sid = str(uuid.uuid4())
        splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
        docs_for_embeddings = []
        metadatas_for_embeddings = []
        
        for t, m in zip(all_texts, all_metas):
            pieces = splitter.split_text(t)
            for i, p in enumerate(pieces):
                docs_for_embeddings.append(p)
                md = m.copy()
                md["chunk"] = i
                metadatas_for_embeddings.append(md)
        
        try:
            # Vector store
            vectorstore = FAISS.from_texts(docs_for_embeddings, embeddings, metadatas=metadatas_for_embeddings)
            storage_path = STORAGE_ROOT / sid
            storage_path.mkdir(parents=True, exist_ok=True)
            vectorstore.save_local(str(storage_path))

            # --- Hybrid: Create BM25 Retriever ---
            bm25_retriever = SimpleBM25Retriever(docs_for_embeddings, metadatas_for_embeddings)

            # Store BOTH in memory
            vectorstore_dict[sid] = vectorstore
            qa_chain_dict[sid] = {"bm25": bm25_retriever}
            qa_history_dict[sid] = []
            
            # Save session to DB
            new_session = ChatSession(
                id=sid,
                user_id=user_id,
                title=f"Session - {datetime.now().isoformat()}"
            )
            db.add(new_session)
            db.commit()
            
            return {
                "session_id": sid,
                "message": f"✅ {len(files)} file(s) uploaded and {len(docs_for_embeddings)} chunks indexed.",
                "errors": errors,
                "action": "created_new"
            }
            
        except Exception as e:
            return JSONResponse({"error": f"Error creating vector store: {str(e)}"}, status_code=500)


class AskRequest(BaseModel):
    question: str
    session_id: str
    session_token: str = None

@app.post("/ask")
async def ask(payload: AskRequest, db: Session = Depends(get_db)):
    session_id = payload.session_id.strip()
    question = payload.question.strip()
    session_token = payload.session_token
    if not session_id:
        return JSONResponse({"error": "Session ID required"}, status_code=400)
    if not question:
        return JSONResponse({"error": "No question provided"}, status_code=400)
    user_id = None
    if session_token:
        user_session = db.query(UserSession).filter(UserSession.session_token == session_token).first()
        if user_session and user_session.expires_at > datetime.now():
            user_id = user_session.user_id

    # --- Ensure vectorstore is loaded ---
    if session_id not in vectorstore_dict:
        storage_path = STORAGE_ROOT / session_id
        if not storage_path.exists():
            return JSONResponse({"error": "Invalid session ID"}, status_code=400)
        try:
            vectorstore = FAISS.load_local(str(storage_path), embeddings, allow_dangerous_deserialization=True)
            vectorstore_dict[session_id] = vectorstore
        except Exception as e:
            return JSONResponse({"error": f"Failed to reload session: {str(e)}"}, status_code=500)
    vectorstore = vectorstore_dict[session_id]

    # --- Get relevant chunks via similarity search ---
    try:
        relevant_docs = vectorstore.similarity_search(question, k=10)  # Increased k to capture more context
    except Exception as e:
        print(f"Similarity search failed: {e}")
        relevant_docs = []

    # --- Enhanced context building with image awareness ---
    image_chunks, code_chunks, text_chunks, table_chunks = [], [], [], []
    
    for doc in relevant_docs:
        meta = getattr(doc, "metadata", {}) or {}
        pc = getattr(doc, "page_content", "")
        chunk_type = meta.get("type", "text")
        
        if chunk_type == "ocr":
            image_chunks.append((meta, pc))
        elif chunk_type == "code":
            code_chunks.append((meta, pc))
        elif chunk_type == "table":
            table_chunks.append((meta, pc))
        elif chunk_type == "text":
            text_chunks.append((meta, pc))

    def format_chunk(meta, pc, is_image=False):
        src = meta.get("source", meta.get("filename", "Unknown"))
        loc = ""
        if "page" in meta:
            loc = f" (Page {meta['page']})"
        elif "paragraph" in meta:
            loc = f" (Paragraph {meta['paragraph']})"
        elif "slide" in meta:
            loc = f" (Slide {meta['slide']})"
        
        # Special handling for image content
        if is_image:
            return f"EXTRACTED TEXT FROM IMAGE: {src}{loc}\n{pc}"
        else:
            return f"Source: {src}{loc}\n{pc[:1500]}"

    # Prioritize image chunks when question is about images
    context_pieces = []
    
    # If question mentions images, prioritize image chunks
    image_keywords = ['image', 'picture', 'photo', 'screenshot', 'diagram', 'chart', 'graph', 'png', 'jpg', 'jpeg']
    has_image_keyword = any(keyword in question.lower() for keyword in image_keywords)
    
    if has_image_keyword and image_chunks:
        context_pieces.extend([format_chunk(meta, pc, is_image=True) for meta, pc in image_chunks])
    
    # Add other content types
    context_pieces.extend([format_chunk(meta, pc) for meta, pc in code_chunks])
    context_pieces.extend([format_chunk(meta, pc) for meta, pc in text_chunks])
    context_pieces.extend([format_chunk(meta, pc) for meta, pc in table_chunks])
    
    # If no image keyword but we have image chunks, include them at the end
    if not has_image_keyword and image_chunks:
        context_pieces.extend([format_chunk(meta, pc, is_image=True) for meta, pc in image_chunks])

    context_text = "\n\n---\n\n".join(context_pieces) or ""

        # COMBINED SMART PROMPT - Handles exact copy + images + all file types
    def build_smart_prompt(question: str, context: str) -> str:
        """Smart prompt that handles exact copy requests, images, and all file types"""
        
        # Exact copy keywords
        exact_keywords = [
            "exactly as given in the file", "exact code", "identical", "verbatim", 
            "copy exactly", "as is", "without changes", "original code",
            "match exactly", "precise copy", "exact same", "unchanged"
        ]
        
        wants_exact = any(keyword in question.lower() for keyword in exact_keywords)
        
        if wants_exact:
            return f"""COPY THE EXACT CONTENT FROM THE UPLOADED FILES:

CONTEXT FROM UPLOADED FILES (includes text, code, tables, and OCR from images):
{context}

USER REQUEST: {question}

YOUR TASK: Copy the content EXACTLY as it appears in the context above.

SPECIAL NOTES:
- "EXTRACTED TEXT FROM IMAGE" means text was obtained via OCR from image files
- Tables are shown in Markdown format
- Code blocks are preserved as-is

STRICT RULES FOR EXACT COPY:
- COPY CHARACTER FOR CHARACTER from context
- NO changes, improvements, or optimizations
- NO explanations or comments
- PRESERVE all formatting, spacing, and comments
- If exact content exists, copy it directly
- If exact content not found, say: "The exact content you requested was not found in the uploaded files."

EXACT COPY:"""
        
        else:
            return f"""Use the context from uploaded files to answer helpfully:

CONTEXT FROM UPLOADED FILES (includes text, code, tables, and OCR from images):
{context}

USER QUESTION: {question}

SPECIAL NOTES:
- "EXTRACTED TEXT FROM IMAGE" means text was obtained via OCR from image files
- Tables are shown in Markdown format - present them as Markdown when relevant
- Code blocks can be improved and optimized

HELPFUL RESPONSE RULES:
- You can improve, optimize, or explain the code
- Add best practices and suggestions
- Fix errors and improve code quality
- Present tables in Markdown format when relevant
- When answering about images, reference that text was extracted via OCR
- If context doesn't have relevant info, provide general help
- Always cite sources when using content from context

HELPFUL RESPONSE:"""

    final_prompt = build_smart_prompt(question, context_text)
    answer = None
    citations = []
    
    if context_text.strip():
        print("#[CONTEXT LENGTH]", len(context_text))
        print("#[CONTEXT SAMPLE]", context_text[:1000])
        print("#[PROMPT SAMPLE]", final_prompt[:1000])
        answer = call_openrouter_model("openai/gpt-oss-20b:free", final_prompt)
        if "not found" in answer.lower():
            answer = None
    
    if not answer or not context_text.strip():
        fallback_prompt = f"Answer the following question as best as possible (no context available, try your best):\nQ: {question}"
        answer = call_openrouter_model("nvidia/nemotron-nano-9b-v2:free", fallback_prompt)
        citations = []
    
    # --- Build citations with image awareness ---
    if context_text.strip() and relevant_docs:
        for doc in relevant_docs:
            md = getattr(doc, "metadata", {}) or {}
            page_content = getattr(doc, "page_content", "")
            
            # Determine location
            location = "Unknown"
            if "page" in md:
                location = f"Page {md['page']}"
            elif "paragraph" in md:
                location = f"Paragraph {md.get('paragraph')}"
            elif "slide" in md:
                location = f"Slide {md.get('slide')}"
            
            # Special handling for image sources
            source_type = md.get("type", "text")
            if source_type == "ocr":
                location = f"Image - {location}" if location != "Unknown" else "Image"
            
            preview = page_content[:200]
            citations.append({
                "source": md.get("source", md.get("filename", "Unknown")),
                "location": location,
                "preview": preview,
                "full_text": page_content,
                "type": md.get("type"),
                "table_data": md.get("table_data")
            })

    # Save messages to database
    new_msg = ChatMessage(session_id=session_id, sender="user", message_text=question)
    db.add(new_msg)
    db.commit()
    ai_msg = ChatMessage(session_id=session_id, sender="ai", message_text=answer)
    db.add(ai_msg)
    db.commit()
    qa_history_dict.setdefault(session_id, []).append({"q": question, "a": answer, "citations": citations})
    
    async def streamer():
        for line in answer.split("\n"):
            yield line + "\n"
            await asyncio.sleep(0.01)
        yield "\n\n---__CITATIONS__---\n"
        yield json.dumps(citations)
    return StreamingResponse(streamer(), media_type="text/plain")


# All your INFO, SAVE, EXPORT, REGISTER, LOGIN, HISTORY, and SESSION endpoints below remain unchanged from your last script.
# (Paste your full original code below this point if this block is not long enough for your entire file).

# ---------- DOCUMENT ANALYSIS ----------
@app.get("/document-analysis/{session_id}")
async def document_analysis(session_id:str):
    try:
        storage_path = STORAGE_ROOT / session_id
        if not storage_path.exists(): return JSONResponse({"error":"Session not found"},status_code=404)
        vectorstore = FAISS.load_local(str(storage_path), embeddings, allow_dangerous_deserialization=True)
        all_docs=[]
        try:
            for i in range(vectorstore.index.ntotal):
                doc = None
                try:
                    doc_id = vectorstore.index_to_docstore_id[i]
                    doc = vectorstore.docstore.search(doc_id)
                except:
                    for k,v in getattr(vectorstore,"docstore",{}).items(): doc=v; break
                if doc and hasattr(doc,"metadata"):
                    meta = doc.metadata
                    preview = doc.page_content[:150]+"..." if len(doc.page_content)>150 else doc.page_content
                    all_docs.append({"source":meta.get("source",meta.get("filename","Unknown")),
                                     "location":meta.get("page",meta.get("paragraph",meta.get("slide","Unknown"))),
                                     "content_preview":preview})
        except:
            try:
                for k,doc in vectorstore.docstore._dict.items():
                    if hasattr(doc,"metadata"):
                        meta = doc.metadata
                        preview = doc.page_content[:150]+"..." if len(doc.page_content)>150 else doc.page_content
                        all_docs.append({"source":meta.get("source",meta.get("filename","Unknown")),
                                         "location":meta.get("page",meta.get("paragraph",meta.get("slide","Unknown"))),
                                         "content_preview":preview})
            except: pass
        sources=defaultdict(list)
        for d in all_docs: sources[d["source"]].append(d)
        return JSONResponse({"document_count":len(all_docs),"source_count":len(sources),"sources":dict(sources)})
    except Exception as e:
        return JSONResponse({"error":str(e)},status_code=500)

# ---------- SAVE CHAT ----------
@app.post("/save-chat")
async def save_chat(request: Request, db: Session = Depends(get_db)):
    body = await request.json()
    session_id = body.get("session_id")
    messages = body.get("messages", [])
    if not session_id or not messages:
        return JSONResponse({"error":"session_id and messages required"},status_code=400)
    saved_count = 0
    for m in messages:
        sender = m.get("sender","user")
        message_text = m.get("message_text","")
        if not message_text.strip(): continue
        new_msg = ChatMessage(session_id=session_id, sender=sender, message_text=message_text)
        db.add(new_msg)
        saved_count += 1
    db.commit()
    qa_history_dict.setdefault(session_id, []).extend(messages)
    return {"status":"ok", "saved": saved_count}

# ---------- EXPORT CHAT ----------
@app.get("/export-chat/{session_id}")
async def export_chat(session_id: str, db: Session = Depends(get_db)):
    msgs = db.query(ChatMessage).filter(ChatMessage.session_id == session_id).order_by(ChatMessage.created_at).all()
    if not msgs:
        return JSONResponse({"error":"No messages found"},status_code=404)
    export_txt = "\n".join([f"{m.sender}: {m.message_text}" for m in msgs])
    return StreamingResponse(BytesIO(export_txt.encode("utf-8")),
                             media_type="text/plain",
                             headers={"Content-Disposition":f"attachment; filename=chat_{session_id}.txt"})

# ---------- GET CHAT HISTORY ----------
@app.get("/chat-history/{session_id}")
async def chat_history(session_id: str):
    history = qa_history_dict.get(session_id, [])
    return {"session_id": session_id, "history": history}

# ---------- AUTH UTILITIES ----------
def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode(), bcrypt.gensalt()).decode()
def verify_password(password: str, hashed: str) -> bool:
    return bcrypt.checkpw(password.encode(), hashed.encode())
def create_session_token() -> str:
    return str(uuid.uuid4())

# ---------- REGISTER ----------
class RegisterRequest(BaseModel):
    username: str
    password: str
    email: str

@app.post("/register")
async def register(payload: RegisterRequest, db: Session = Depends(get_db)):
    username = payload.username.strip()
    password = payload.password.strip()
    email = payload.email.strip()
    if not username or not password or not email:
        return JSONResponse({"error":"All fields required"}, status_code=400)
    existing_user = db.query(User).filter(User.username == username).first()
    if existing_user:
        return JSONResponse({"error":"Username already exists"}, status_code=400)
    import hashlib
    password_hash = hashlib.sha256(password.encode()).hexdigest()
    new_user = User(username=username, password_hash=password_hash, email=email)
    db.add(new_user)
    db.commit()
    return {"message":"User registered successfully", "username": username}

# ---------- LOGIN ----------
class LoginRequest(BaseModel):
    username: str
    password: str

@app.post("/login")
async def login(payload: LoginRequest, db: Session = Depends(get_db)):
    username = payload.username.strip()
    password = payload.password.strip()
    if not username or not password:
        return JSONResponse({"error":"Username and password required"}, status_code=400)
    user = db.query(User).filter(User.username == username).first()
    if not user:
        return JSONResponse({"error":"Invalid username or password"}, status_code=401)
    import hashlib
    password_hash = hashlib.sha256(password.encode()).hexdigest()
    if user.password_hash != password_hash:
        return JSONResponse({"error":"Invalid username or password"}, status_code=401)
    token = str(uuid.uuid4())
    expires_at = datetime.now() + timedelta(hours=24)
    user_session = UserSession(user_id=user.id, session_token=token, expires_at=expires_at)
    db.add(user_session)
    db.commit()
    return {"message":"Login successful", "session_token": token, "user_id": user.id, "username": username}

# ---------- GET CHAT MESSAGES FOR A SESSION ----------
@app.get("/chat-messages/{session_id}")
async def get_chat_messages(session_id: str, db: Session = Depends(get_db)):
    session = db.query(ChatSession).filter(ChatSession.id == session_id).first()
    if not session:
        return JSONResponse({"error": "Session not found"}, status_code=404)
    messages = [
        {
            "sender": msg.sender,
            "messagetext": msg.message_text,
            "createdat": msg.created_at.isoformat(),
            "citations": msg.citations if msg.citations else []
        }
        for msg in session.messages
    ]
    return {"session_id": session_id, "messages": messages}

@app.get("/user-sessions/{user_id}")
async def get_user_sessions(
    user_id: str,
    authorization: str = Header(None),
    db: Session = Depends(get_db)
):
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing or invalid Authorization header")
    session_token = authorization.split(" ")[1]
    user_session = db.query(UserSession).filter(
        UserSession.user_id == user_id,
        UserSession.session_token == session_token,
        UserSession.expires_at > datetime.now()
    ).first()
    if not user_session:
        raise HTTPException(status_code=401, detail="Invalid or expired session token")
    sessions = db.query(ChatSession).filter(ChatSession.user_id == user_id).order_by(ChatSession.created_at.desc()).all()
    result = []
    for s in sessions:
        last_msg = db.query(ChatMessage).filter(ChatMessage.session_id == s.id).order_by(ChatMessage.created_at.desc()).first()
        preview = last_msg.message_text[:100] + "..." if last_msg else ""
        history = []
        msgs = db.query(ChatMessage).filter(ChatMessage.session_id == s.id).order_by(ChatMessage.created_at).all()
        for m in msgs:
            history.append({"sender": m.sender, "message_text": m.message_text})
        result.append({
            "session_id": s.id,
            "title": s.title or "Untitled Session",
            "saved_at": s.created_at.isoformat(),
            "history": history,
            "preview": preview
        })
    return {"sessions": result}

@app.post("/download-table")
async def download_table(request: Request):
    data = await request.json()
    table_data = data.get("table_data", [])
    filename = data.get("source", "table.csv")
    if not table_data or not isinstance(table_data, list):
        return JSONResponse({"error": "No table data found"}, status_code=400)
    import csv, io
    output = io.StringIO()
    writer = csv.writer(output)
    for row in table_data:
        writer.writerow(row)
    output.seek(0)
    # Use `filename` if you want custom filenames!
    return StreamingResponse(output, media_type="text/csv",
                            headers={"Content-Disposition":f"attachment; filename={filename if filename.endswith('.csv') else filename+'.csv'}"})
